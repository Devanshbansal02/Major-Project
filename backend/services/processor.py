import io
import logging
import re
from pathlib import Path

import fitz  # PyMuPDF
import pytesseract
from PIL import Image, ImageEnhance
from pptx import Presentation

from backend.config import PROCESSED_DIR, RAW_DIR

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Text Preprocessing & Extraction
# ---------------------------------------------------------------------------


def _clean_extracted_text(text: str) -> str:
    """Clean up raw extracted text (especially from OCR) before embedding."""
    if not text:
        return ""

    # Normalize carriage returns
    text = text.replace("\r\n", "\n").replace("\r", "\n")

    lines = []
    for line in text.split("\n"):
        line = line.strip()
        # Remove purely non-alphanumeric/noise lines (e.g., '---', '• • •', etc.)
        # Keep lines that have at least some word characters
        if line and re.search(r"\w", line):
            # Normalize multiple spaces inside the line
            line = re.sub(r"[ \t]+", " ", line)
            lines.append(line)

    # Re-join lines. We use double newline to preserve paragraph semantics for chunking.
    cleaned = "\n\n".join(lines)

    # Remove instances of 3 or more consecutive newlines
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)

    return cleaned.strip()


def _pptx_to_text(path: Path) -> str:
    prs = Presentation(str(path))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def _pdf_typed_to_text(path: Path) -> str:
    doc = fitz.open(str(path))
    parts = [page.get_text() for page in doc]
    doc.close()
    return "\n".join(parts)


def _preprocess_image(img: Image.Image) -> Image.Image:
    # Convert to grayscale
    img = img.convert("L")
    # Increase contrast
    enhancer = ImageEnhance.Contrast(img)
    img = enhancer.enhance(2.0)
    return img


def _pdf_handwritten_to_text(path: Path) -> str:
    doc = fitz.open(str(path))
    parts = []
    for page in doc:
        pix = page.get_pixmap(dpi=300)
        img = Image.open(io.BytesIO(pix.tobytes("png")))
        img = _preprocess_image(img)
        parts.append(pytesseract.image_to_string(img))
    doc.close()
    return "\n".join(parts)


def _docx_to_text(path: Path) -> str:
    try:
        from docx import Document as DocxDocument  # python-docx
    except ImportError:
        logger.error("python-docx not installed: cannot extract DOCX")
        return ""
    doc = DocxDocument(str(path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    # Also pull text from tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text = cell.text.strip()
                if text:
                    paragraphs.append(text)
    return "\n".join(paragraphs)


def _image_to_text(path: Path) -> str:
    img = Image.open(str(path))
    img = _preprocess_image(img)
    return pytesseract.image_to_string(img)


# ---------------------------------------------------------------------------
# File metadata helper
# ---------------------------------------------------------------------------


def get_file_metadata(path: Path) -> dict:
    ext = path.suffix.lower()
    meta: dict = {"file_size_bytes": path.stat().st_size, "extension": ext}
    try:
        if ext == ".pdf":
            doc = fitz.open(str(path))
            meta["page_count"] = len(doc)
            meta["word_count"] = sum(len(p.get_text().split()) for p in doc)
            doc.close()
        elif ext == ".pptx":
            prs = Presentation(str(path))
            meta["slide_count"] = len(prs.slides)
    except Exception as exc:
        logger.warning("Could not read metadata for %s: %s", path, exc)
    return meta


# ---------------------------------------------------------------------------
# Core: process a DB note row → extracted text
# ---------------------------------------------------------------------------


async def process_note_file(note: dict) -> tuple[Path | None, str]:
    """Extract text from an uploaded note.

    Args:
        note: DB row dict with keys: id, filename, file_type

    Returns:
        (output_path, text): output_path is the cached .txt file.
        Returns (None, "") on failure.
    """
    note_id = note["id"]
    filename = note["filename"]
    ftype = note["file_type"]

    raw_path = Path(RAW_DIR) / filename
    out_path = Path(PROCESSED_DIR) / f"{note_id}.txt"

    # Return cached extraction if available
    if out_path.exists():
        text = out_path.read_text(encoding="utf-8", errors="replace")
        return out_path, text

    if not raw_path.exists():
        logger.error("Raw file missing for note_id=%s: %s", note_id, raw_path)
        return None, ""

    try:
        if ftype == "pptx":
            text = _pptx_to_text(raw_path)
        elif ftype == "pdf_handwritten":
            text = _pdf_handwritten_to_text(raw_path)
        elif ftype == "pdf_typed":
            text = _pdf_typed_to_text(raw_path)
        elif ftype == "docx":
            text = _docx_to_text(raw_path)
        elif ftype == "image":
            text = _image_to_text(raw_path)
        else:
            logger.warning("Unknown file type %s for note_id=%s", ftype, note_id)
            text = f"[Unsupported file type: {ftype}]"

        # Apply cleaning preprocessing
        text = _clean_extracted_text(text)

        out_path.parent.mkdir(parents=True, exist_ok=True)
        out_path.write_text(text, encoding="utf-8")
        logger.info(
            "Extracted & cleaned text for note_id=%s (%d words)",
            note_id,
            len(text.split()),
        )
    except Exception as exc:
        logger.error("Extraction failed note_id=%s: %s", note_id, exc, exc_info=True)
        return None, ""

    return out_path, text


# ---------------------------------------------------------------------------
# Bulk ingestion: extract + embed into ChromaDB, update is_embedded flag
# ---------------------------------------------------------------------------


async def ingest_note_files(notes: list[dict]) -> dict:
    """Process and embed a list of DB note rows.

    Updates the `is_embedded` flag in the DB for each successfully indexed note.
    """
    from backend.db import get_db
    from backend.services import rag_service

    Path(PROCESSED_DIR).mkdir(parents=True, exist_ok=True)

    processed = 0
    for note in notes:
        note_id = note["id"]
        subject_id = note["subject_id"]

        _, text = await process_note_file(note)
        if not text.strip():
            logger.warning("Empty text for note_id=%s: skipping", note_id)
            continue

        chunks_added = await rag_service.add_note_to_index(subject_id, note_id, text)
        if chunks_added > 0:
            processed += 1
            # Mark as embedded in DB
            try:
                db = await get_db()
                await db.execute(
                    "UPDATE notes SET is_embedded = 1 WHERE id = ?", (note_id,)
                )
                await db.commit()
                await db.close()
            except Exception as exc:
                logger.error(
                    "Failed to update is_embedded for note_id=%s: %s", note_id, exc
                )

    logger.info("Ingestion complete: %d/%d notes embedded", processed, len(notes))
    return {"status": "done", "processed": processed, "total": len(notes)}
